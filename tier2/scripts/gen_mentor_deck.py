#!/usr/bin/env python3
"""
gen_mentor_deck.py — Generate mentor-meeting presentation for Tier 2 epoll UAF research.

Outputs:
  tier2/docs/TIER2_MENTOR_DECK.pptx        (primary, 22 slides + speaker notes)
  tier2/docs/TIER2_MENTOR_EXEC_SUMMARY.pdf (2-page backup via weasyprint)

Data sources (parsed manually into tables below):
  tier2/docs/EXPERIMENT_INDEX.md, VERIFICATION_LEDGER.md, DEAD_ENDS_REGISTER.md,
  ASSUMPTIONS_REGISTER.md, RUNNER_GUIDE.md, tier2/evidence/*_RESULTS.md
"""

import html
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path("/mnt/work/company/cyphermatrix/repos/bad-epoll-lab")
OUT_PPTX = REPO / "tier2" / "docs" / "TIER2_MENTOR_DECK.pptx"
OUT_PDF = REPO / "tier2" / "docs" / "TIER2_MENTOR_EXEC_SUMMARY.pdf"

# ----------------------------------------------------------------------------
# Design constants
# ----------------------------------------------------------------------------
NAVY = RGBColor(0x0F, 0x24, 0x40)
NAVY2 = RGBColor(0x1B, 0x3A, 0x5C)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
AMBER = RGBColor(0xF5, 0x7F, 0x17)
BLUE = RGBColor(0x15, 0x65, 0xC0)
TEXT = RGBColor(0x21, 0x21, 0x21)
MUTED = RGBColor(0x5F, 0x6B, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF9, 0xFC)
LIGHT_BORDER = RGBColor(0xD8, 0xDE, 0xE6)

BODY = "Calibri"
MONO = "Consolas"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

_slide_no = [0]


def new_slide(bg=BG):
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    _slide_no[0] += 1
    return slide


def add_text(slide, x, y, w, h, lines, wrap=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    for i, (txt, size, bold, color, font, space_after) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font or BODY
    return tb


def rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def rounded(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        sh.adjustments[0] = 0.045
    except Exception:
        pass
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def header(slide, kicker, title, accent=AMBER):
    rect(slide, 0, 0, SW, Inches(1.02), NAVY)
    rect(slide, 0, Inches(1.02), SW, Pt(3.2), accent)
    add_text(slide, Inches(0.55), Inches(0.10), Inches(12.2), Inches(0.30),
             [(kicker, 12, False, RGBColor(0x9F, 0xB4, 0xCC), None, 0)])
    add_text(slide, Inches(0.55), Inches(0.36), Inches(12.2), Inches(0.62),
             [(title, 26, True, WHITE, None, 0)])
    add_text(slide, Inches(0.55), Inches(1.18), Inches(12.3), Inches(0.30),
             [("bad-epoll-lab  ·  CVE-2026-46242  ·  tier2-android-port  ·  Intern progress review", 9, False, MUTED, None, 0)])


def footer(slide, note=None):
    n = _slide_no[0]
    line = f"{n}"
    add_text(slide, Inches(12.55), Inches(7.14), Inches(0.6), Inches(0.3),
             [(line, 10, True, MUTED, None, 0)], align=PP_ALIGN.RIGHT)
    if note:
        add_text(slide, Inches(0.55), Inches(7.14), Inches(11.8), Inches(0.3),
                 [(note, 8, False, MUTED, None, 0)])


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def chip(slide, x, y, w, h, label, color, size=11, bold=True):
    sh = rounded(slide, x, y, w, h, color)
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = WHITE
    r.font.name = BODY
    return sh


def bullets(slide, x, y, w, h, items, size=13, gap=5):
    """items: list of (text, level, color, bold) or plain str."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(item, str):
            txt, lvl, color, bold, fsize = item, 0, TEXT, False, size
        else:
            txt, lvl, color, bold, fsize = item[0], item[1], item[2], item[3], (item[4] if len(item) > 4 else size)
        prefix = "\u25CF  " if lvl == 0 else "\u2013  "
        r = p.add_run()
        r.text = prefix + txt
        r.font.size = Pt(fsize)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = BODY
    return tb


def _norm(lines, size, color):
    """Accept (text, size, color, bold[, space_after]) or plain str; return add_text 6-tuples."""
    out = []
    for item in lines:
        if isinstance(item, str):
            out.append((item, size, False, color, None, 0))
        elif len(item) == 4:
            txt, sz, col, b = item
            out.append((txt, sz, b, col, None, 0))
        elif len(item) == 5:
            txt, sz, col, b, sp = item
            out.append((txt, sz, b, col, None, sp))
        else:
            out.append(tuple(item))
    return out


def card(slide, x, y, w, h, title, body, accent=BLUE, title_size=14, body_size=11):
    rounded(slide, x, y, w, h, WHITE, LIGHT_BORDER)
    rect(slide, x, y + Inches(0.06), Inches(0.07), h - Inches(0.12), accent)
    add_text(slide, x + Inches(0.22), y + Inches(0.10), w - Inches(0.4), Inches(0.35),
             [(title, title_size, True, accent, None, 0)])
    add_text(slide, x + Inches(0.22), y + Inches(0.46), w - Inches(0.4), h - Inches(0.55),
             _norm(body, body_size, TEXT), wrap=True)


def make_table(slide, x, y, w, h, headers, rows, col_w, font_size=9.5,
               row_h=Inches(0.24), header_fill=NAVY, align_map=None):
    nrows, ncols = len(rows) + 1, len(headers)
    gf = slide.shapes.add_table(nrows, ncols, x, y, w, h)
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    for ci, cw in enumerate(col_w):
        tbl.columns[ci].width = cw
    for ri in range(nrows):
        tbl.rows[ri].height = row_h
    for ci, htxt in enumerate(headers):
        c = tbl.cell(0, ci)
        c.fill.solid()
        c.fill.fore_color.rgb = header_fill
        tf = c.text_frame
        tf.margin_left = tf.margin_right = Pt(4)
        tf.margin_top = tf.margin_bottom = Pt(1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if (align_map and ci in align_map) else PP_ALIGN.LEFT
        r = p.add_run()
        r.text = htxt
        r.font.size = Pt(font_size)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = BODY
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(0xEE, 0xF2, 0xF7)
            tf = c.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Pt(4)
            tf.margin_top = tf.margin_bottom = Pt(1)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if (align_map and ci in align_map) else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = val
            r.font.size = Pt(font_size)
            r.font.bold = False
            r.font.color.rgb = TEXT
            r.font.name = BODY if not val.startswith("0x") else MONO
    return tbl


# ----------------------------------------------------------------------------
# SLIDE 1 — Title
# ----------------------------------------------------------------------------
s = new_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.62), SW, Pt(3), AMBER)
add_text(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.4),
         [("TIER 2 EXPLOITABILITY ASSESSMENT  —  PROGRESS REVIEW", 14, True, RGBColor(0x9F, 0xB4, 0xCC), None, 0)])
add_text(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(1.7),
         [("CVE-2026-46242 (\"Bad Epoll\")", 40, True, WHITE, None, 0),
          ("Use-After-Free in the Linux epoll subsystem on Android ARM64 GKI", 20, False, RGBColor(0xC9, 0xD6, 0xE8), None, 0)])
add_text(s, Inches(0.9), Inches(3.75), Inches(11.5), Inches(0.7),
         [("Phase: Tier 2 — exploitability assessment on linux-6.12.67 (Android 14 GKI, commit 7e35917775b8)", 14, False, WHITE, None, 2),
          ("Branch: tier2-android-port  ·  HEAD cc0dc7754 (verified on origin via git ls-remote)", 14, False, WHITE, None, 0)])
add_text(s, Inches(0.9), Inches(4.95), Inches(11.5), Inches(0.8),
         [("Intern: Aayush Bankar  ·  Mentor meeting  ·  August 2026", 16, True, WHITE, None, 4),
          ("Agenda: What we proved  ·  What we disproved  ·  The one blocker  ·  Where to go next", 13, False, RGBColor(0xC9, 0xD6, 0xE8), None, 0)])
notes(s, "Title slide. Context: Tier 1 (x86_64, Fedora, kernelCTF COS kernel) already produced a full LPE chain "
         "(UAF->cross-cache->AAR->KASLR->ROP->commit_creds, UID 0) — see article/CVE-2026-46242_Tier1_Final_Writeup.md. "
         "Tier 2 moves the same CVE to Android ARM64 GKI where PAC/BTI/MTE/SELinux apply and the exploit chain must be rebuilt from scratch.")

# ----------------------------------------------------------------------------
# SLIDE 2 — One-slide verdict
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "EXECUTIVE SUMMARY", "One-Slide Verdict")
card(s, Inches(0.55), Inches(1.6), Inches(6.0), Inches(1.85), "STATUS  —  AMBER",
     [("The UAF is real and 100% reproducible — but only with debugger assistance.", 12, TEXT, False, 2),
      ("Natural (no-GDB) race: 0 / 102,740 attempts across two campaigns.", 12, TEXT, False, 0)], AMBER)
card(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(1.85), "PROVED  —  GREEN",
     [("Race mechanics captured by hardware watchpoint (VER-026, EXP-015).", 12, TEXT, False, 2),
      ("msg_msg reclaim of freed struct eventpoll in kmalloc-192 (VER-027, EXP-018).", 12, TEXT, False, 0)], GREEN)
card(s, Inches(0.55), Inches(3.6), Inches(6.0), Inches(1.85), "DISPROVED  —  RED",
     [("All 4 exploitation chains dead; 21 dead ends documented.", 12, TEXT, False, 2),
      ("Dual-watch KASLR leak retracted (VER-033): conditions mutually exclusive.", 12, TEXT, False, 0)], RED)
card(s, Inches(6.8), Inches(3.6), Inches(6.0), Inches(1.85), "PRIMITIVE  —  BLUE",
     [("Only UAF write: NULL at offset 160 of a freed kmalloc-192 object.", 12, TEXT, False, 2),
      ("EXP-016 audit: fixed NULL value + fixed offset => DoS-only on this config.", 12, TEXT, False, 0)], BLUE)
rounded(s, Inches(0.55), Inches(5.62), Inches(12.25), Inches(1.25), NAVY)
add_text(s, Inches(0.85), Inches(5.78), Inches(11.7), Inches(1.0),
         [("Recommendation: timebox 2 weeks on physical ARM64 hardware (timing-widening, Path A) — then conclude DoS-only (Path C) if 0 hits.", 14, True, WHITE, None, 3),
          ("Ask: approve the hardware timebox, or conclude DoS-only now?", 12, False, RGBColor(0xC9, 0xD6, 0xE8), None, 0)])
footer(s)
notes(s, "Verdict slide — deliver in 2 minutes. Every claim is evidence-backed: VER-026 (EXP-015 HW watchpoint trace of the "
         "outer-close/inner-close race writing NULL at offset 160 of the FREED inner_epoll), VER-027 (EXP-018 msg_msg reclaim). "
         "Chains dead: VER-028, VER-031, VER-033. Natural schedulability: VER-034 (0/10k) and VER-039 (0/92,740). "
         "The one open question is hardware timing variance, which QEMU TCG cannot model.")

# ----------------------------------------------------------------------------
# SLIDE 3 — What we proved
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "FOUNDATION", "What We Proved (GDB-Assisted, Hardware-Watchpoint Verified)", accent=GREEN)
bullets(s, Inches(0.55), Inches(1.55), Inches(12.3), Inches(3.6), [
    ("CVE-2026-46242 = close-vs-close race: Thread A closes the OUTER epoll, Thread B closes the INNER epoll fd that the outer watches.", 0, TEXT, False),
    ("VER-026 (EXP-015, hardware watchpoint trace): Thread A executes WRITE_ONCE(file->f_ep, NULL) in __ep_remove, then Thread B's __fput sees f_ep == NULL, bypasses eventpoll_release_file via the lockless fast path (include/linux/eventpoll.h), and frees the inner struct eventpoll.", 0, TEXT, False),
    ("Thread A resumes and hlist_del_rcu(&epi->fllink) writes NULL at offset 160 (refs.first) of the FREED inner eventpoll — the UAF write, captured by a hardware watchpoint.", 0, TEXT, False),
    ("VER-020: sizeof(struct eventpoll) = 176 B => kmalloc-192 (generic cache).", 0, TEXT, False),
    ("VER-027 (EXP-018/019): msg_msg with 144 B user payload reclaims the freed slot; attacker controls every byte from slab offset 48 onward (marker 0xdead000000000000 landed exactly at offset 136).", 0, TEXT, False),
    ("Methodology: every claim is labeled STATIC (source/disassembly) or RUNTIME (traced execution) — never conflated (protocol Rule 6).", 0, TEXT, False),
], size=13)
rounded(s, Inches(0.55), Inches(5.35), Inches(12.25), Inches(1.6), WHITE, LIGHT_BORDER)
add_text(s, Inches(0.85), Inches(5.5), Inches(11.7), Inches(1.35),
         [("Evidence quote (EXP-015_unified_trace.log):", 11, True, MUTED, None, 2),
          ("\"Thread B hit ep_free(...)! Thread B finished ep_free. target_epoll is now FREED!\" / \"Hardware watchpoint 6 ... Old value = -281474920794800  New value = 0  __hlist_del (n=...) at ./include/linux/list.h:989\"", 11, False, TEXT, MONO, 0)])
footer(s)
notes(s, "Key message: the vulnerability mechanism is fully understood and proven at runtime, not just from source. "
         "Quote the two lines from EXP-015 evidence verbatim. The lockless fast path is the crux: eventpoll_release() "
         "does a lockless READ_ONCE(file->f_ep) and returns early if NULL — that is what lets Thread B free the object "
         "out from under Thread A. Tier 1 on x86_64 already proved full LPE for the same CVE (see article writeup PDF).")

# ----------------------------------------------------------------------------
# SLIDE 4 — What we disproved
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "DEAD ENDS", "What We Disproved (4 chains, 21 documented dead ends)", accent=RED)
bullets(s, Inches(0.55), Inches(1.55), Inches(12.3), Inches(5.2), [
    ("Chain 0 — controlled crash via percpu_counter_dec: operates on the OUTER (valid) epoll, never the freed INNER one. VER-028 (EXP-019).", 0, TEXT, False),
    ("Chain 1 — dual-watch KASLR leak: single-epitem UAF and multi-epitem pointer write are mutually exclusive. VER-033 (EXP-024) retracts VER-029/030.", 0, TEXT, False),
    ("Chain 2 — arbitrary decrement via fake user_struct: ep parameter in __ep_remove is always the OUTER epoll; GDB redirect of outer_epoll->user failed, fbc = root_user+8. VER-031/032 (EXP-023b).", 0, TEXT, False),
    ("Chain 3 — full LPE (modprobe_path / creds): depends on Chains 0-2 => blocked.", 0, TEXT, False),
    ("struct epitem same-cache reclaim: list_del_init corrupts offsets 24/32 BEFORE the free; INIT_LIST_HEAD of a fresh epitem overwrites the corruption. VER-016 (EXP-008).", 0, TEXT, False),
    ("struct file UAF: filp cache is SLAB_TYPESAFE_BY_RCU (unmergeable); stale f_op/f_lock accesses are benign. VER-018/025 (EXP-010/015).", 0, TEXT, False),
    ("21 dead ends total (DE-001..DE-021): 4 chains, 5 target objects, 4 primitive expansions, 4 spray/reclaim paths, 4 early misconceptions.", 0, TEXT, True),
    ("Protocol Rule 8: retracted claims stay visible in the ledger (VER-010, VER-029, VER-030) with reasons and corrections.", 0, TEXT, False),
], size=13)
footer(s)
notes(s, "Walk the mentor through the dead-ends register (tier2/docs/DEAD_ENDS_REGISTER.md). Emphasise that these are not "
         "beliefs — each is killed by a hardware-watchpoint trace or source-level structural proof (VER IDs cited). "
         "The most intellectually interesting one is Chain 1 (VER-033): you need single-epitem for f_ep=NULL (the UAF), but "
         "you need 2+ epitems for the write to be a non-NULL kernel pointer. No topology satisfies both — structural impossibility.")

# ----------------------------------------------------------------------------
# SLIDE 5 — The blocker
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "BLOCKER", "The Blocker: Natural Race Is Not Winnable (Root Cause)", accent=AMBER)
bullets(s, Inches(0.55), Inches(1.55), Inches(12.3), Inches(3.9), [
    ("All UAF hits so far are DEBUGGER-ASSISTED: GDB patches Thread A into an infinite loop inside __ep_remove, creating an artificial preemption window (assumptions A-039/A-040 falsified).", 0, TEXT, False),
    ("cond_resched() at fs/eventpoll.c lines 888 & 903 is a NO-OP in the 2-CPU pinned PREEMPT_VOLUNTARY setup (NAT-002 CORRECTION, VER-035):", 0, TEXT, False),
    ("   dynamic_cond_resched static key is FALSE (CONFIG_PREEMPT_DYNAMIC, default voluntary) => returns 0 immediately.", 1, TEXT, False),
    ("   TIF_NEED_RESCHED is never set on CPU 0: timer tick with queued=0 never calls resched_curr(); no cross-CPU IPIs in the pinned idle topology.", 1, TEXT, False),
    ("__ep_remove itself contains ZERO preemption points (verified in disassembly) — the whole vulnerable sequence runs atomically with respect to scheduling.", 0, TEXT, False),
    ("The real race window is pure instruction-count timing: ~250-550 cycles (~125-275 ns @ 2 GHz) — below scheduling granularity.", 0, TEXT, False),
    ("QEMU TCG emulation does not model hardware memory-bus / cache-coherency behaviour, so emulated timing diverges from real silicon.", 0, TEXT, False),
], size=13)
card(s, Inches(0.55), Inches(5.5), Inches(12.25), Inches(1.45), "Consequence",
     [("Statistical result: 0/102,740 natural attempts (NAT-001: 10,000 simple; NAT-005: 92,740 with isolcpus=1, nohz_full, rcu_nocbs and a 4 MB cache-eviction sweeper). Best alignment error: 1 cycle (~16 ns).", 12, TEXT, False, 0)], AMBER)
footer(s)
notes(s, "This is the heart of the meeting. Prior work (Tier 1, x86_64) won the race because the COS kernel/exploit used "
         "different scheduling conditions. On this Android GKI build, the two conditions that make the race winnable "
         "(a scheduling yield inside the window, and hardware timing variance) are both absent: PREEMPT_VOLUNTARY + pinned "
         "idle CPUs => cond_resched no-op; TCG => no real bus/cache timing. Therefore all 'successes' were debugger-forced.")

# ----------------------------------------------------------------------------
# SLIDE 6 — Evidence dashboard
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "DATA", "Evidence Dashboard", accent=BLUE)
rows = [
    ["Experiments executed", "19 (EXP-006..024, NAT-001/002/005, AND-001)", "EXPERIMENT_INDEX.md"],
    ["GDB-assisted UAF hits", "~100% on demand (EXP-015/018/019/024)", "EXP-015_unified_trace.log"],
    ["Natural race hits", "0 / 102,740 (NAT-001: 10k + NAT-005: 92,740)", "NAT-001, NAT-005 logs"],
    ["Best timing alignment error", "1 cycle (~16 ns) @ delay=2,360 cycles", "NAT-005_raw_serial.log"],
    ["Verified claims (active)", "VER-009 ... VER-039 (20 entries)", "VERIFICATION_LEDGER.md"],
    ["Retracted claims", "3 (VER-010, VER-029, VER-030) — kept visible per Rule 8", "VERIFICATION_LEDGER.md"],
    ["Dead ends", "21 (4 chains · 5 objects · 4 primitives · 4 sprays · 4 misconceptions)", "DEAD_ENDS_REGISTER.md"],
    ["Android portability", "AND-001 PASSED — SysV IPC functional, load_msg trapped", "AND-001_raw_ipc.log"],
    ["Kernel under test", "linux-6.12.67, commit 7e35917775b8, CONFIG_SYSVIPC=y, PREEMPT_VOLUNTARY", ".config"],
]
make_table(s, Inches(0.55), Inches(1.55), Inches(12.25), Inches(5.1), ["Metric", "Value", "Evidence file"],
           rows, [Inches(3.1), Inches(6.3), Inches(2.85)], font_size=11, row_h=Inches(0.52))
footer(s)
notes(s, "Reference table — do not read aloud, point to it. All numbers traceable to committed raw logs in tier2/evidence/. "
         "Protocol Rule 1: evidence lives in the repo, never in agent session state. Rule 2: no VERIFIED without reading and "
         "quoting the raw log — this is how the earlier fabricated claims (VER-001/002/004/006/007) were caught and purged.")

# ----------------------------------------------------------------------------
# SLIDE 7 — Vulnerability deep dive
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "TECHNICAL DEEP-DIVE", "Vulnerability Mechanics: Outer-Close / Inner-Close Race")
bullets(s, Inches(0.55), Inches(1.55), Inches(6.4), Inches(5.3), [
    ("Topology: epoll_ctl(EPOLL_CTL_ADD, outer, inner_fd) — one epitem; inner_epoll->refs gains epi->fllink.", 0, TEXT, False),
    ("Thread A: close(outer_fd) -> ep_clear_and_put -> __ep_remove(outer_ep, epi_for_inner).", 0, TEXT, False),
    ("eventpoll.c:826 — single-epitem check (head->first == &epi->fllink && !epi->fllink.next)", 0, TEXT, False),
    ("eventpoll.c:828 — WRITE_ONCE(file->f_ep, NULL)  [file = inner epoll file]", 0, TEXT, False),
    ("Thread B: close(inner_fd) -> __fput -> eventpoll_release(): lockless READ_ONCE(f_ep) == NULL => bypasses eventpoll_release_file => ep_free => kfree(inner_epoll) to kmalloc-192.", 0, TEXT, False),
    ("Thread A resumes -> eventpoll.c:836 hlist_del_rcu(&epi->fllink): UAF NULL write at freed+160 (refs.first).", 0, TEXT, True),
    ("Lines 840-857 (rb_erase_cached, spin_lock_irq, percpu_counter_dec) all operate on the OUTER epoll — valid object.", 0, TEXT, False),
], size=12)
card(s, Inches(7.2), Inches(1.55), Inches(5.6), Inches(2.5), "struct eventpoll (176 B -> kmalloc-192)",
     [("lock (spinlock) ............ @ 96", 10.5, TEXT, False, 0),
      ("rbr (rb_root_cached) ....... @ 104-119", 10.5, TEXT, False, 0),
      ("user (user_struct *) ....... @ 136", 10.5, TEXT, False, 0),
      ("refs (hlist_head, f_ep) .... @ 160  <== UAF write target", 10.5, RED, True, 0)], BLUE)
card(s, Inches(7.2), Inches(4.2), Inches(5.6), Inches(2.6), "msg_msg spray mapping",
     [("msg_msg header = 48 B; user payload starts at slab offset 48", 10.5, TEXT, False, 2),
      ("slab offset X  =  user payload byte (X - 48)", 10.5, TEXT, False, 2),
      ("offset 96 -> byte 48 | offset 136 -> byte 88 | offset 160 -> byte 112", 10.5, TEXT, False, 2),
      ("144 B payload => 192 B total => same kmalloc-192 cache", 10.5, TEXT, False, 0)], BLUE)
footer(s)
notes(s, "Full mechanics: show the sequence on the left, layouts on the right. The ONLY UAF operation on the freed object is "
         "hlist_del_rcu at offset 160 (VER-032). Everything else in __ep_remove touches the outer epoll which remains alive — "
         "this single fact killed Chains 0 and 2. The refs field (hlist_head) is used for f_ep / eventpoll_release_file bookkeeping.")

# ----------------------------------------------------------------------------
# SLIDE 8 — Chain 0
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "CHAIN 0 — DEAD (VER-028)", "Controlled Crash via percpu_counter_dec", accent=RED)
card(s, Inches(0.55), Inches(1.55), Inches(6.0), Inches(1.7), "Theory",
     [("Spray the freed inner_epoll with msg_msg placing 0xDEAD000000000000 at offset 136 (ep->user); Thread A's percpu_counter_dec would dereference garbage -> controlled panic.", 11, TEXT, False, 0)], BLUE)
card(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(1.7), "Result",
     [("Reclaim works (marker landed exactly at offset 136) — but NO crash. Thread A completed __ep_remove cleanly.", 11, TEXT, False, 0)], RED)
card(s, Inches(0.55), Inches(3.4), Inches(6.0), Inches(1.7), "Why it failed",
     [("percpu_counter_dec(&ep->user->epoll_watches) uses the ep parameter = OUTER eventpoll (valid), not the freed INNER one. The offset-136 bytes are never read in this race path.", 11, TEXT, False, 0)], RED)
card(s, Inches(6.8), Inches(3.4), Inches(6.0), Inches(1.7), "Lesson",
     [("Attacker control of a field means nothing if the code path never reads that field. The only field the freed object's bytes feed is offset 160 via hlist_del_rcu.", 11, TEXT, False, 0)], AMBER)
rounded(s, Inches(0.55), Inches(5.35), Inches(12.25), Inches(1.5), WHITE, LIGHT_BORDER)
add_text(s, Inches(0.85), Inches(5.5), Inches(11.7), Inches(1.2),
         [("Evidence (EXP-019_raw_gdb.log): \"Marker at 0xffff0000035c4a48 (offset 136): 0xdead000000000000 ... MATCH\" then \"No crash occurred - hang timeout (25s) triggered\".", 11, False, TEXT, MONO, 0)])
footer(s)
notes(s, "Chain 0 was the 'controlled crash PoC' — a valid CVE trigger demo. It failed for a subtle reason: in the "
         "outer-close/inner-close race, the ep parameter of __ep_remove is always the outer epoll. So even though the spray "
         "put our poison at offset 136 of the freed slot, nothing dereferences it. Positive spin: the msg_msg reclaim primitive "
         "is rock-solid (VER-027). Full writeup: tier2/evidence/EXP-019_RESULTS.md.")

# ----------------------------------------------------------------------------
# SLIDE 9 — Chain 1
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "CHAIN 1 — RETRACTED (VER-033)", "Dual-Watch KASLR Leak via hlist_del_rcu", accent=RED)
card(s, Inches(0.55), Inches(1.55), Inches(6.0), Inches(1.7), "Theory (EXP-022b)",
     [("Add inner_epoll to TWO outer epolls => 2 epitems in inner_epoll->refs. Closing outer1 makes hlist_del_rcu write epi2's address (a kernel pointer) at offset 160 -> read back via msgrcv -> KASLR defeat.", 11, TEXT, False, 0)], BLUE)
card(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(1.7), "Original evidence was flawed",
     [("EXP-022b log: 'AFTER spray' dump captured BEFORE ep_free; labels inverted; log ends in GDB error with no msgrcv readback ever captured.", 11, TEXT, False, 0)], AMBER)
card(s, Inches(0.55), Inches(3.4), Inches(12.25), Inches(1.7), "EXP-024 re-test with hardware watchpoint — result",
     [("With 2+ epitems, the single-epitem check (eventpoll.c:826) is FALSE => WRITE_ONCE(f_ep, NULL) never executes => the lockless bypass (VER-026) cannot trigger => inner_epoll is NEVER freed before hlist_del_rcu. Watchpoint fired on LIVE memory (ep_free_inner_seen: False), value NULL.", 11, TEXT, False, 0)], RED)
card(s, Inches(0.55), Inches(5.25), Inches(12.25), Inches(1.6), "Structural impossibility (VER-033)",
     [("UAF requires single-epitem (f_ep=NULL); kernel-pointer write requires multi-epitem. No topology satisfies both. VER-029/VER-030 retracted and kept visible.", 11, TEXT, True, 0)], RED)
footer(s)
notes(s, "This is the retraction story — valuable for the mentor to see the self-correction process working (protocol Rules 2, 8). "
         "EXP-024 used a hardware watchpoint on inner_epoll+160 to settle it definitively: the write went to LIVE memory. "
         "Quote: \"VERDICT: hlist_del_rcu wrote to LIVE inner_epoll (never freed during trace).\" "
         "Evidence: tier2/evidence/EXP-024_RESULTS.md.")

# ----------------------------------------------------------------------------
# SLIDE 10 — Chain 2
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "CHAIN 2 — DEAD (VER-031/032)", "Arbitrary Decrement via Fake user_struct", accent=RED)
card(s, Inches(0.55), Inches(1.55), Inches(6.0), Inches(1.7), "Theory (EXP-023b)",
     [("Second msg_msg spray holds a fake struct user_struct; point ep->user (offset 136) at it so percpu_counter_dec decrements an attacker-chosen address (modprobe_path).", 11, TEXT, False, 0)], BLUE)
card(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(1.7), "What we did",
     [("GDB wrote the fake user_struct into kernel memory and overwrote outer_epoll->user (offset 136) to point at it; breakpoint at percpu_counter_add_batch.", 11, TEXT, False, 0)], BLUE)
card(s, Inches(0.55), Inches(3.4), Inches(12.25), Inches(1.7), "Result",
     [("Breakpoint hit: fbc = 0xffff80008152bea0 = root_user + 8. The kernel read outer_epoll->user as the REAL root_user, not our fake pointer. modprobe_path unchanged (still \"/sbin/modprobe\"). Redirect failed.", 11, TEXT, False, 0)], RED)
card(s, Inches(0.55), Inches(5.25), Inches(12.25), Inches(1.6), "Why",
     [("Same root cause as Chain 0: percpu_counter_dec always uses the OUTER epoll (valid, never freed, never reclaimable). The only UAF write remains hlist_del_rcu at offset 160.", 11, TEXT, True, 0)], RED)
footer(s)
notes(s, "Evidence: tier2/evidence/EXP-023b_RESULTS.md and EXP-023b_raw_gdb.log. GDB actively tried to help the chain and it "
         "still failed — the redirect is structurally impossible in this race because the decrement runs on the outer epoll. "
         "The RUNNER_GUIDE's Chain-2 design assumed ep->user of the FREED eventpoll would be used; that assumption was wrong.")

# ----------------------------------------------------------------------------
# SLIDE 11 — Chain 3 + primitive analysis
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "CHAIN 3 — BLOCKED + EXP-016 AUDIT", "Full LPE and the Remaining Primitive", accent=RED)
card(s, Inches(0.55), Inches(1.55), Inches(6.0), Inches(2.0), "Chain 3: Full LPE",
     [("Depends on Chain 1 (KASLR leak) + Chain 2 (arbitrary write) — both dead.", 11, TEXT, False, 2),
      ("Remaining primitive: fixed-value NULL write at fixed offset 160 of a reclaimed kmalloc-192 slot.", 11, TEXT, False, 0)], RED)
card(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(2.0), "EXP-016: kmalloc-192 offset-160 audit",
     [("All generic-kmalloc 176-192 B structs checked via pahole against the real vmlinux + source audit.", 11, TEXT, False, 2),
      ("Verdict: NULLing offset 160 yields crash-only or benign effects in every candidate.", 11, TEXT, True, 0)], BLUE)
make_table(s, Inches(0.55), Inches(3.75), Inches(12.25), Inches(3.1),
           ["Candidate", "Field @160", "NULL effect", "Verdict"],
           [
               ["fib6_info", "rcu.func", "NULL RCU callback call", "Crash only (DoS)"],
               ["snd_timer_user", "ioctl_lock.wait_list.next", "mutex slowpath NULL deref", "Crash only (DoS)"],
               ["packet_fanout", "arr[3] (sock*)", "fanout_demux NULL sock", "Crash only (DoS)"],
               ["urb", "interval (int)", "interval=0", "Benign"],
               ["wakeup_source", "expire_count", "counter=0", "Benign"],
           ],
           [Inches(2.3), Inches(3.3), Inches(4.1), Inches(2.55)], font_size=10.5, row_h=Inches(0.5))
footer(s)
notes(s, "EXP-016 is the 'clean negative' that defines the exploit ceiling: with a NULL-only write at offset 160, no "
         "reachable kmalloc-192 object converts that into a controlled read/write/call. fib6_info (via unshare+netlink "
         "RTM_NEWROUTE) gives the most reliable DoS. This is why the verdict is 'DoS-only on this configuration' unless "
         "natural reachability changes or a different race variant is found. Full audit: tier2/docs/EXP-016_RESULTS.md.")

# ----------------------------------------------------------------------------
# SLIDE 12 — Phase 3
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "PHASE 3 — NATURAL SCHEDULABILITY", "Can the Race Be Won Without GDB?", accent=BLUE)
bullets(s, Inches(0.55), Inches(1.55), Inches(12.3), Inches(3.1), [
    ("NAT-001 (10,000 iterations, CPU pinning + SCHED_FIFO + timing widening): 0 / 10,000 hits — FAILED (VER-034).", 0, TEXT, False),
    ("NAT-002: preemption-point audit — cond_resched at lines 888/903 exists but is a NO-OP (see next slide). CORRECTED (VER-035).", 0, TEXT, False),
    ("NAT-005 (92,740 iterations, closed-loop adaptive search):", 0, TEXT, True),
    ("   Kernel cmdline: isolcpus=1 nohz_full=1 rcu_nocbs=1 (CPU 1 isolated from timer ticks / RCU)", 1, TEXT, False),
    ("   4 MB cache-eviction sweeper worker hammering L1/L2 to induce memory-bus refill delays", 1, TEXT, False),
    ("   Empirically calibrated window: 2,330 cycles; best alignment: 2,360 cycles with 1-cycle (~16 ns) error", 1, TEXT, False),
    ("   Result: 0 / 92,740 UAF hits — PASSED as a negative result (VER-039)", 1, RED, True),
    ("Caveat: QEMU TCG does not model hardware cache-coherency / memory-bus timing; results reflect emulated basic-block scheduling.", 0, AMBER, True),
], size=13)
card(s, Inches(0.55), Inches(5.35), Inches(12.25), Inches(1.5), "Still PLANNED",
     [("NAT-003 (msg_msg reclaim under natural timing) · NAT-004 (per-CPU partial-slab interference) · AND-002 (KASLR on/off) · AND-003 (SELinux audit) · AND-004 (MTE/KASAN_HW_TAGS).", 11, TEXT, False, 0)], BLUE)
footer(s)
notes(s, "Phase 3 was designed as a falsification campaign (EVO-009): all prior evidence was debugger-assisted, so we built "
         "statistical tests to answer 'is the race naturally reachable?'. Answer so far: no, in QEMU. The 1-cycle near-miss "
         "at delay=2,360 cycles shows the harness calibration is essentially perfect — yet zero hits, which strongly suggests "
         "the window cannot be hit by deterministic instruction timing alone under TCG. Evidence: NAT-005_RESULTS.md + raw log.")

# ----------------------------------------------------------------------------
# SLIDE 13 — NAT-002 correction
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "WHY: NAT-002 CORRECTION", "cond_resched() Is a No-Op Here — the Mechanism", accent=AMBER)
bullets(s, Inches(0.55), Inches(1.55), Inches(12.3), Inches(4.1), [
    ("cond_resched() -> dynamic_cond_resched() (static_call): CONFIG_PREEMPT_DYNAMIC=y, default PREEMPT_VOLUNTARY => static key is FALSE => returns 0 immediately without yielding.", 0, TEXT, False),
    ("Even if it called __cond_resched(): should_resched() requires preempt_count()==0 AND TIF_NEED_RESCHED set.", 0, TEXT, False),
    ("TIF_NEED_RESCHED is NEVER set on our CPU 0: timer tick with queued=0 does not call resched_curr(); Thread B on CPU 1 sends no scheduler IPIs (pinned, idle topology).", 0, TEXT, False),
    ("Disassembly confirms: bl dynamic_cond_resched at lines 888/903 of ep_clear_and_put; __ep_remove has ZERO cond_resched/might_resched inside.", 0, TEXT, False),
    ("Consequence: the multi-epitem topology only buys an instruction-count window (~250-550 cycles) between processing epi_N and epi_{N+1} — no scheduling yield.", 0, TEXT, True),
    ("Escalation ideas if pursued: runtime PREEMPT_DYNAMIC switch via debugfs (requires /sys/kernel/debug), a real PREEMPT kernel build, or hardware with real timer variance.", 0, AMBER, False),
], size=13)
footer(s)
notes(s, "The NAT-002 correction is a great example of falsification: the initial audit claimed the cond_resched points were "
         "real preemption windows (STATIC-HIGH-CONFIDENCE), then mechanism analysis proved they cannot yield in our topology. "
         "Files: tier2/evidence/NAT-002/NAT-002_RESULTS.md (initial) and NAT-002_CORRECTION.md (correction). "
         "This is also the argument for Path A: on real hardware, timer interrupts, IPIs, and memory-bus contention DO create "
         "the variance that TCG cannot emulate.")

# ----------------------------------------------------------------------------
# SLIDE 14 — Android portability
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "ANDROID PORTABILITY", "AND-001: SysV IPC Works on the Target Kernel", accent=GREEN)
card(s, Inches(0.55), Inches(1.55), Inches(6.0), Inches(2.1), "AND-001 — PASSED (VER-038)",
     [("msgget / msgsnd / msgrcv functional; kernel call trapped at load_msg (ipc/msgutil.c:96, addr 0xffff800080430100) during static-musl harness execution.", 11, TEXT, False, 2),
      ("120 B payload + 48 B header = 168 B => kmalloc-192 => msg_msg spray remains viable for eventpoll reclaim on Android.", 11, TEXT, True, 0)], GREEN)
card(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(2.1), "Context caveats",
     [("Test ran in the QEMU initramfs as root (shell context), NOT inside an app under SELinux enforcing + seccomp.", 11, TEXT, False, 2),
      ("AND-003 (SELinux syscall policy audit) and app-context seccomp checks are still PLANNED.", 11, AMBER, True, 0)], AMBER)
bullets(s, Inches(0.55), Inches(3.85), Inches(12.3), Inches(3.0), [
    ("Kernel: same source commit 7e35917775b8 (Android 14 GKI 6.12.67) used for the QEMU build — portability of the primitive is direct.", 0, TEXT, False),
    ("PAC / BTI: never reached — the surviving primitive is a data-only NULL write, so PAC/BTI are irrelevant unless a new chain is found (assumption A-028, hypothesis).", 0, TEXT, False),
    ("MTE / KASAN_HW_TAGS: AND-004 PLANNED — MTE on production devices may turn the UAF into a detected fault (DoS) rather than silent corruption.", 0, AMBER, False),
    ("KASLR: AND-002 PLANNED — KASLR does not change race reachability (offsets are resolved per-boot), only leak-based chains would care.", 0, TEXT, False),
    ("SELinux: even with UID 0, Android confines root via SELinux contexts — Tier 1 writeup already flagged that LPE on Android would need selinux_state patching.", 0, TEXT, False),
], size=13)
footer(s)
notes(s, "AND-001 closes the 'does SysV IPC even exist on Android GKI' question — it does, CONFIG_SYSVIPC=y. This was a "
         "CRITICAL BLOCKER in the assumptions register (A-024/A-025) and is now resolved for the shell context. The remaining "
         "Android unknowns are app-context seccomp/SELinux (AND-003) and MTE (AND-004). Evidence: tier2/evidence/AND-001_RESULTS.md.")

# ----------------------------------------------------------------------------
# SLIDE 15 — Three paths forward
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "DECISION", "Three Paths Forward")
make_table(s, Inches(0.55), Inches(1.55), Inches(12.25), Inches(3.9),
           ["", "Path A: Hardware timing-widening", "Path B: Alternative race variant", "Path C: Conclude DoS-only"],
           [
               ["Idea", "Physical ARM64 device; false-sharing cache-line bouncing, slab contention, IPI/timer storms", "Find a race where the freed object IS the ep parameter (e.g., different epoll topology / other subsystems)", "Document all 21 dead ends + statistical negatives; publish negative-result writeup"],
               ["Effort", "2-4 weeks", "2-4 weeks research", "~1 week"],
               ["Risk", "Medium (hardware access needed)", "High (structural dead ends likely)", "Low"],
               ["Success prob.", "Medium — only path with theoretical basis; TCG cannot model the variable", "Low — source shows zero preemption inside __ep_remove", "High confidence in conclusion"],
               ["Outcome if ok", "Natural UAF hit -> resume chain work", "New primitive -> resume chain work", "Defensible, honest negative result"],
           ],
           [Inches(1.5), Inches(3.8), Inches(3.8), Inches(3.15)], font_size=10, row_h=Inches(0.7))
rounded(s, Inches(0.55), Inches(5.75), Inches(12.25), Inches(1.1), NAVY)
add_text(s, Inches(0.85), Inches(5.9), Inches(11.7), Inches(0.85),
         [("Hybrid recommendation: Path A for a strict 2-week timebox with kill-criteria; if 0 hits in ~1M iterations with all widening techniques -> Path C.", 13, True, WHITE, None, 0)])
footer(s)
notes(s, "Decision-support slide. Present all three honestly. Path A is the only one with a theoretical basis (real hardware "
         "adds the timing variance that TCG lacks). Path B is speculative — the source shows no preemption inside __ep_remove, "
         "so a different race would need a different freeing thread (e.g., fd table teardown ordering). Path C is always "
         "available and is the professionally correct fallback: a documented negative result for a real CVE assessment.")

# ----------------------------------------------------------------------------
# SLIDE 16 — Recommendation & ask
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "RECOMMENDATION", "What I Propose + What I Need", accent=GREEN)
bullets(s, Inches(0.55), Inches(1.55), Inches(12.3), Inches(3.2), [
    ("Primary: run Path A as a 2-week timebox on physical ARM64 hardware (or an ARM64 KVM host).", 0, TEXT, True),
    ("Techniques: false-sharing cache-line bouncing (3rd thread), slab/grooming contention, IPI storms, timerfd storms, and closed-loop delay calibration ported from NAT-005.", 0, TEXT, False),
    ("Kill-criteria: 0 natural UAF hits across ~1M iterations with all widening techniques enabled => conclude Path C.", 0, TEXT, False),
    ("Fallback deliverable: final negative-result writeup (exploitability assessment) covering the CVE, all chains, and the DoS-only conclusion — suitable as a research note / CVE assessment.", 0, TEXT, False),
    ("Whatever the outcome, evidence protocol continues: every run logged in EXPERIMENT_INDEX, VER ledger updated, raw logs committed and pushed.", 0, TEXT, False),
], size=13)
card(s, Inches(0.55), Inches(5.1), Inches(12.25), Inches(1.7), "What I need from you (mentor)",
     [("1) Decision: approve the 2-week hardware timebox (Path A), or conclude DoS-only now (Path C)?", 12, TEXT, False, 3),
      ("2) Access: ARM64 test device (e.g., unlocked dev unit) or ARM64 KVM host + ability to read kernel logs.", 12, TEXT, False, 0)], GREEN)
footer(s)
notes(s, "End the meeting with the ask. If the mentor says 'conclude', we execute Path C: final report + PDF writeup. "
         "If 'timebox', we start hardware work immediately. Either way the repo evidence discipline continues unchanged.")

# ----------------------------------------------------------------------------
# SLIDE 17 — Risk register
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "RISKS", "Risk Register")
rows = [
    ["TCG vs real hardware timing divergence", "Medium", "Mitigation: Path A on physical device; treat QEMU results as lower-bound evidence"],
    ["SELinux/seccomp may block msgget/msgsnd in real app context", "Medium", "AND-003 planned; AND-001 tested shell context only"],
    ["Scheduler divergence (2-CPU QEMU vs big.LITTLE device)", "Low-Med", "Retest timing on device before concluding"],
    ["Device SLUB config may differ from QEMU .config", "Low", "Fetch device /proc/slabinfo before spray work"],
    ["MTE on production devices detects the UAF -> fault, not corruption", "Unknown", "AND-004 planned; may cap impact at DoS regardless"],
    ["Timebox creep on hardware tuning", "Medium", "Hard 2-week cap with kill-criteria (slide 16)"],
    ["False confidence from GDB-assisted success", "Handled", "Phase 3 falsification campaign (EVO-009); labels STATIC/RUNTIME everywhere"],
    ["PAC/BTI if a new control-flow chain is found", "Low now", "Data-only primitive avoids them; re-assess only if new chain appears"],
]
make_table(s, Inches(0.55), Inches(1.55), Inches(12.25), Inches(4.9), ["Risk", "Likelihood", "Mitigation / note"],
           rows, [Inches(4.6), Inches(1.5), Inches(6.15)], font_size=10.5, row_h=Inches(0.55))
footer(s)
notes(s, "Honest risk inventory for the meeting. The 'false confidence' row is handled by design: Phase 3 exists precisely "
         "to falsify debugger-assisted successes. MTE is the sleeper risk: on production Android with MTE enabled, the UAF "
         "may be caught at hlist_del_rcu time — turning any future exploit into a crash. That caps long-term exploitability "
         "on modern devices regardless of race reachability.")

# ----------------------------------------------------------------------------
# SLIDE 18 — Appendix: experiment index
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "APPENDIX", "Experiment Index (19 executed + planned)")
rows = [
    ["EXP-006", "07-23", "pipe_buffer cross-cache struct-file UAF", "INCONCLUSIVE"],
    ["EXP-007", "07-27", "same-cache epitem reclaim", "PASSED"],
    ["EXP-008", "07-28", "timing: list_del_init before free", "PASSED (dead end)"],
    ["EXP-009", "07-28", "struct file reclaim via open() spray", "PASSED"],
    ["EXP-010", "07-29", "pipe_buffer content control", "PASSED (dead end)"],
    ["EXP-014", "07-30", "eventpoll_release_file stale-epi path", "DISPROVED"],
    ["EXP-015", "07-31", "HW trace: lockless-path race (VER-026)", "VERIFIED"],
    ["EXP-016", "07-31", "kmalloc-192 offset-160 target audit", "CLEAN NEGATIVE"],
    ["EXP-018", "08-01", "msg_msg reclaim (VER-027)", "PASSED"],
    ["EXP-019", "08-01", "Chain 0 controlled crash", "FAILED (VER-028)"],
    ["EXP-022b", "08-01", "Chain 1 dual-watch leak", "RETRACTED (VER-033)"],
    ["EXP-023b", "08-01", "Chain 2 arbitrary decrement", "NEGATIVE (VER-031)"],
    ["EXP-024", "08-02", "dual-watch re-test", "PASSED negative (VER-033)"],
    ["NAT-001", "08-02", "10,000-iter natural race", "FAILED 0/10,000"],
    ["NAT-002", "08-02", "preemption-point audit", "CORRECTED"],
    ["NAT-005", "08-05", "92,740-iter closed-loop search", "PASSED negative (VER-039)"],
    ["AND-001", "08-05", "SysV IPC on target kernel", "PASSED (VER-038)"],
]
make_table(s, Inches(0.55), Inches(1.55), Inches(12.25), Inches(5.2),
           ["ID", "Date", "Objective", "Result"], rows,
           [Inches(1.1), Inches(1.0), Inches(7.0), Inches(3.15)], font_size=9.5, row_h=Inches(0.3))
footer(s)
notes(s, "EXP-011/012/013 also ran (type-confusion unreachable VER-025; wait-queue drain identity VER-022; single-thread "
         "write-before-free VER-021) — see VERIFICATION_LEDGER.md. Full registry: tier2/docs/EXPERIMENT_INDEX.md. "
         "PLANNED: NAT-003, NAT-004, AND-002, AND-003, AND-004.")

# ----------------------------------------------------------------------------
# SLIDE 19 — Appendix: verification ledger highlights
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "APPENDIX", "Verification Ledger Highlights (VER-014 .. VER-039)")
rows = [
    ["VER-014", "eventpoll_epi is isolated dedicated cache (SLAB_ACCOUNT) — no merge into kmalloc-128", "STATIC"],
    ["VER-015", "same-cache epitem reclaim occupies exact freed slot", "RUNTIME"],
    ["VER-016", "epitem same-cache reclaim dead: rdllink corrupted before free", "RUNTIME"],
    ["VER-018", "filp cache SLAB_TYPESAFE_BY_RCU — struct file cross-cache dead", "STATIC"],
    ["VER-020", "struct eventpoll = 176 B -> kmalloc-192; refs @ offset 160", "STATIC"],
    ["VER-025", "struct file UAF benign: f_op/f_lock accesses safe (filp-only cache)", "STATIC"],
    ["VER-026", "RACE PROVEN: HW watchpoint, NULL write @ freed+160 via lockless path", "RUNTIME"],
    ["VER-027", "msg_msg 144 B reclaims freed eventpoll; attacker control 48+", "RUNTIME"],
    ["VER-028", "Chain 0 dead: percpu_counter_dec uses OUTER epoll", "RUNTIME"],
    ["VER-031", "Chain 2 dead: decrement always root_user (outer epoll)", "RUNTIME"],
    ["VER-032", "Only UAF write = hlist_del_rcu @ offset 160 (NULL, single-epitem)", "RUNTIME"],
    ["VER-033", "Chain 1 retracted: UAF and pointer-write mutually exclusive", "RUNTIME"],
    ["VER-034", "NAT-001: 0/10,000 natural hits (CI upper bound 0.0384%)", "RUNTIME"],
    ["VER-035", "NAT-002 corrected: cond_resched no-op in 2-CPU pinned voluntary", "STATIC"],
    ["VER-038", "AND-001: SysV IPC functional; load_msg trapped", "RUNTIME"],
    ["VER-039", "NAT-005: 0/92,740; best alignment error 1 cycle (~16 ns)", "RUNTIME"],
]
make_table(s, Inches(0.55), Inches(1.55), Inches(12.25), Inches(5.2),
           ["VER", "Claim", "Method"], rows,
           [Inches(1.1), Inches(9.7), Inches(1.45)], font_size=9.5, row_h=Inches(0.3))
footer(s)
notes(s, "Complete ledger: tier2/docs/VERIFICATION_LEDGER.md. Retractions VER-010/029/030 kept visible per Rule 8. "
         "Evolution notes EVO-005 (dedicated epi cache correction) and EVO-009 (Phase 3 falsification pivot) are the two "
         "biggest understanding shifts of the project.")

# ----------------------------------------------------------------------------
# SLIDE 20 — Appendix: dead ends
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "APPENDIX", "Dead Ends Register — 21 Entries (DE-001 .. DE-021)")
rows = [
    ["DE-001..004", "Chains 0-3 (crash / leak / decrement / LPE)", "VER-028, VER-033, VER-031"],
    ["DE-005", "struct file UAF via ep_item_poll type confusion", "VER-025 (EXP-015)"],
    ["DE-006", "struct file UAF via pipe_buffer cross-cache", "VER-018 (EXP-010)"],
    ["DE-007", "epitem same-cache reclaim (EPOLL_CTL_ADD)", "VER-016 (EXP-008)"],
    ["DE-008", "snd_timer_user reclaim of epitem/eventpoll", "EVO-005"],
    ["DE-009", "snd_timer_user mutex slowpath corruption", "VER-013"],
    ["DE-010", "multi-epitem kernel-pointer write @160", "VER-033 (EXP-024)"],
    ["DE-011", "rb_erase_cached arbitrary write", "VER-032 (EXP-023b)"],
    ["DE-012", "spin_lock_irq lock corruption @96", "VER-032 (EXP-023b)"],
    ["DE-013", "KASLR defeat via dual-watch leak", "VER-030 retracted"],
    ["DE-014", "pipe_buffer spray for struct file", "VER-018 (EXP-010)"],
    ["DE-015", "snd_timer_user spray for epitem", "EVO-005"],
    ["DE-016", "eventpoll_epi merge into kmalloc-128", "VER-014 (STATIC)"],
    ["DE-017", "cross-cache grooming into eventpoll_epi", "VER-014 (STATIC)"],
    ["DE-018..021", "early misconceptions (simple race, timing spikes, level-3 claim, release_file mutex path)", "FAILURE_ANALYSIS.md"],
]
make_table(s, Inches(0.55), Inches(1.55), Inches(12.25), Inches(5.2),
           ["ID", "Dead path", "Killing evidence"], rows,
           [Inches(1.7), Inches(6.9), Inches(3.65)], font_size=10, row_h=Inches(0.31))
footer(s)
notes(s, "Full register: tier2/docs/DEAD_ENDS_REGISTER.md. Revisit allowed only under new conditions: changed kernel config "
         "(e.g., CONFIG_PREEMPT=y), new code discovery, or new hardware capability — each requiring a new VER with RUNTIME "
         "evidence that contradicts the killing evidence. This guards against re-litigating settled questions.")

# ----------------------------------------------------------------------------
# SLIDE 21 — Appendix: commands
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "APPENDIX", "Commands & Reproducibility")
add_text(s, Inches(0.55), Inches(1.55), Inches(12.3), Inches(3.6),
         [("Compile harness (static musl):", 12, True, NAVY, MONO, 2),
          ("./aarch64-linux-musl-cross/bin/aarch64-linux-musl-gcc -static -O0 -g -o rootfs/harness scripts/<H>.c -pthread", 11, False, TEXT, MONO, 6),
          ("Package initramfs:", 12, True, NAVY, MONO, 2),
          ("cd rootfs && chmod +x init harness && find . -print0 | cpio --null -ov --format=newc > ../initramfs.cpio", 11, False, TEXT, MONO, 6),
          ("Boot QEMU with GDB stub:", 12, True, NAVY, MONO, 2),
          ("DEBUG=1 ./scripts/run_qemu.sh > /dev/null 2>&1 &   # gdb port :1234, cmdline: kasan=on nokaslr", 11, False, TEXT, MONO, 6),
          ("Run a GDB automation script:", 12, True, NAVY, MONO, 2),
          ("gdb -batch -q -x scripts/<S>.py android/artifacts/vmlinux", 11, False, TEXT, MONO, 6),
          ("Stop QEMU / verify push:", 12, True, NAVY, MONO, 2),
          ("pkill -f qemu-system-aarch64 ; git ls-remote origin tier2-android-port", 11, False, TEXT, MONO, 0)])
card(s, Inches(0.55), Inches(5.35), Inches(12.25), Inches(1.5), "Key paths",
     [("Kernel: third_party/linux-6.12.67/ (vmlinux, .config)  ·  Evidence: tier2/evidence/  ·  Scripts: tier2/scripts/  ·  Docs: tier2/docs/  ·  Cross-compiler: tier2/aarch64-linux-musl-cross/", 11, TEXT, False, 0)], BLUE)
footer(s)
notes(s, "For the mentor (or any future agent): full protocol in tier2/docs/EXPERIMENT_PROTOCOL.md (10 non-negotiable rules), "
         "runner guide in RUNNER_GUIDE.md. The AGENTS.md at repo root summarises the critical commands and the evidence "
         "standard. Everything is reproducible from the committed sources.")

# ----------------------------------------------------------------------------
# SLIDE 22 — Appendix: git provenance
# ----------------------------------------------------------------------------
s = new_slide()
header(s, "APPENDIX", "Git History & Repository State")
bullets(s, Inches(0.55), Inches(1.55), Inches(6.5), Inches(4.2), [
    ("Branch: tier2-android-port; HEAD cc0dc7754 — \"feat(exp-nat005): implement isolcpus=1 and 4MB cache-eviction race widening\".", 0, TEXT, False),
    ("origin verified: git ls-remote shows cc0dc7754f0bf34e9c8d596877113c9b00d014ca (2026-08-07).", 0, GREEN, True),
    ("Recent commits:", 0, TEXT, True),
    ("  853fd4fe2 — NAT-005 calibration + 2,330-cycle closed-loop search", 1, TEXT, False),
    ("  f904dda78 / 8851d70b3 — closed-loop adaptive search, 100k iterations", 1, TEXT, False),
    ("  b96d5d0ed / 37888baf5 — QEMU CPU topology + launch-ahead harness", 1, TEXT, False),
    ("  8f1392cab — AND-001 SysV IPC verification", 1, TEXT, False),
    ("  1f17b0a71 / 47d9d11c7 — NAT-001 complete: 0/10,000, race not naturally winnable", 1, TEXT, False),
    ("  b40f0ae20 — EXP-024: retract VER-029/030", 1, TEXT, False),
    ("  7c6e67392 / 92733a977 / 4b965dc3b / dcebc159b — EXP-023b/022b/019/018 evidence", 1, TEXT, False),
], size=12)
card(s, Inches(7.3), Inches(1.55), Inches(5.5), Inches(4.2), "Working-tree honesty (protocol Rule 10)",
     [("NOT fully clean as of this deck:", 11, True, RED, False, 2),
      ("  Untracked: AGENTS.md, EXP-014/EXP-020/EXP-021 evidence + scripts (pending a commit)", 10, TEXT, False, 2),
      ("  Modified post-commit: AND-001_raw_ipc.log, EXP-012/018 raw logs, NAT-005_topology_raw.log", 10, TEXT, False, 2),
      ("  third_party/security-research submodule has untracked content", 10, TEXT, False, 2),
      ("Action: one evidence commit (exp-020/021 + stray logs) before final report.", 10, AMBER, True, 0)], AMBER)
footer(s)
notes(s, "Per protocol Rule 10 I am explicitly flagging that the working tree is NOT fully clean: EXP-014/020/021 artifacts "
         "and AGENTS.md are uncommitted, and several raw logs were modified after their commits. Plan: bundle them into one "
         "evidence commit, push, and re-verify with git ls-remote before the final report. This transparency is itself part "
         "of the evidence discipline the mentor should see.")

# ----------------------------------------------------------------------------
# Build companion HTML executive summary for PDF backup
# ----------------------------------------------------------------------------
def build_pdf_backup():
    try:
        import weasyprint  # noqa
    except Exception:
        print("[!] weasyprint unavailable — skipping PDF backup", file=sys.stderr)
        return
    css = """
    @page { size: A4; margin: 2cm 1.8cm; @bottom-right { content: counter(page) " / " counter(pages); font-size: 8pt; color: #888; } }
    body { font-family: 'DejaVu Sans', sans-serif; font-size: 10pt; color: #212121; line-height: 1.45; }
    h1 { font-size: 17pt; color: #0F2440; border-bottom: 3px solid #F57F17; padding-bottom: 4px; }
    h2 { font-size: 12.5pt; color: #0F2440; margin-top: 14px; border-bottom: 1px solid #D8DEE6; padding-bottom: 2px; }
    table { border-collapse: collapse; width: 100%; font-size: 8.6pt; margin: 6px 0; }
    th { background: #0F2440; color: #fff; padding: 3px 6px; text-align: left; }
    td { border: 1px solid #D8DEE6; padding: 3px 6px; }
    tr:nth-child(even) td { background: #EEF2F7; }
    .verdict { background: #FFF8E1; border: 1px solid #F57F17; padding: 8px 10px; }
    .dead { color: #C62828; font-weight: bold; }
    .ok { color: #2E7D32; font-weight: bold; }
    code { font-family: 'DejaVu Sans Mono', monospace; font-size: 8.4pt; background: #F1F3F5; padding: 0 2px; }
    """
    body = html.escape  # shorthand
    tables = []

    def tbl(headers, rows):
        h = "".join(f"<th>{body(x)}</th>" for x in headers)
        r = "".join("<tr>" + "".join(f"<td>{body(c)}</td>" for c in row) + "</tr>" for row in rows)
        return f"<table><tr>{h}</tr>{r}</table>"

    html_doc = f"""<!DOCTYPE html><html><head><meta charset="utf-8"><style>{css}</style></head><body>
<h1>CVE-2026-46242 Tier 2 Progress — Executive Summary</h1>
<p><b>Project:</b> bad-epoll-lab · <b>Branch:</b> tier2-android-port @ cc0dc7754 (verified on origin) · <b>Kernel:</b> linux-6.12.67, Android 14 GKI (commit 7e35917775b8) · <b>Date:</b> 2026-08-07</p>
<div class="verdict"><b>Verdict:</b> The epoll UAF is real and reproducible with debugger assistance (hardware-watchpoint proven, VER-026), and msg_msg reclaim of the freed <code>struct eventpoll</code> in kmalloc-192 is reliable (VER-027). However, the race is <b>not naturally schedulable</b>: 0 / 102,740 unaided attempts (VER-034, VER-039). All four exploitation chains are structurally dead (21 dead ends); the only primitive is a fixed NULL write at offset 160, which EXP-016 shows is DoS-only on this configuration. <b>Recommendation:</b> 2-week hardware timebox (timing-widening on physical ARM64) with kill-criteria, else conclude DoS-only.</div>
<h2>1. What was proved (RUNTIME evidence)</h2>
<ul>
<li>VER-026 (EXP-015): two-threaded race — Thread A clears <code>f_ep</code>, Thread B bypasses <code>eventpoll_release_file</code> via the lockless fast path and frees <code>inner_epoll</code>; Thread A's <code>hlist_del_rcu</code> writes NULL at offset 160 of the freed object. Captured by hardware watchpoint.</li>
<li>VER-027 (EXP-018/019): msg_msg with 144 B payload reclaims the freed slot with attacker control from byte 48.</li>
<li>VER-020: <code>sizeof(struct eventpoll)=176</code> → kmalloc-192. VER-038: SysV IPC works on the target kernel (AND-001).</li>
</ul>
<h2>2. What was disproved (each with killing evidence)</h2>
<ul>
<li class="dead">Chain 0 — percpu_counter_dec crash: uses OUTER epoll, never the freed inner (VER-028).</li>
<li class="dead">Chain 1 — dual-watch KASLR leak: single-epitem UAF and multi-epitem pointer write are mutually exclusive (VER-033; VER-029/030 retracted).</li>
<li class="dead">Chain 2 — arbitrary decrement via fake user_struct: decrement always runs on root_user (VER-031/032).</li>
<li class="dead">Chain 3 — full LPE: depends on Chains 1+2. Only remaining primitive: NULL @ offset 160 → EXP-016 audit shows crash-only/benign outcomes for all reachable kmalloc-192 structs (fib6_info, snd_timer_user, packet_fanout...).</li>
<li class="dead">epitem same-cache reclaim (VER-016), struct file UAF (VER-018/025), snd_timer_user theories (EVO-005/VER-013).</li>
</ul>
<h2>3. Why the natural race fails</h2>
<ul>
<li>cond_resched() at eventpoll.c:888/903 is a NO-OP: dynamic_cond_resched static key is FALSE and TIF_NEED_RESCHED is never set on the pinned, idle CPU (NAT-002 correction, VER-035).</li>
<li>__ep_remove has zero preemption points (disassembly-confirmed). The window is pure instruction timing: ~250–550 cycles (~125–275 ns @ 2 GHz).</li>
<li>QEMU TCG does not model hardware cache-coherency/memory-bus timing. NAT-005 (92,740 iterations, isolcpus=1, nohz_full, 4 MB cache-eviction sweeper) reached a 1-cycle (~16 ns) best alignment error with 0 hits.</li>
</ul>
<h2>4. Where we are</h2>
{tbl(["Metric", "Value"], [
    ["Experiments executed", "19 (EXP-006..024, NAT-001/002/005, AND-001)"],
    ["GDB-assisted UAF hits", "~100% on demand"],
    ["Natural race hits", "0 / 102,740"],
    ["Best timing alignment error", "1 cycle (~16 ns)"],
    ["Dead ends", "21 (4 chains · 5 objects · 4 primitives · 4 sprays · 4 misconceptions)"],
    ["Retracted claims (kept visible)", "VER-010, VER-029, VER-030"],
    ["Active verification entries", "VER-009 .. VER-039"],
    ["Android portability", "AND-001 PASSED (SysV IPC); AND-002/003/004 planned"],
])}
<h2>5. Decision options</h2>
{tbl(["Path", "Description", "Effort", "Risk", "Probability"], [
    ["A: Hardware timing-widening", "Physical ARM64 device; false-sharing cache bouncing, slab contention, IPI/timer storms", "2–4 weeks", "Medium", "Medium — only path with theoretical basis"],
    ["B: Alternative race variant", "Race where the freed object IS the ep parameter of __ep_remove", "2–4 weeks", "High", "Low"],
    ["C: Conclude DoS-only", "Document all dead ends + statistical negatives; publish negative-result writeup", "~1 week", "Low", "High confidence in conclusion"],
])}
<p><b>Hybrid recommendation:</b> Path A for a strict 2-week timebox with kill-criteria (0 hits in ~1M iterations ⇒ Path C).</p>
<h2>6. Reproducibility pointers</h2>
<ul>
<li>Protocol: <code>tier2/docs/EXPERIMENT_PROTOCOL.md</code> (10 rules) · Runner: <code>tier2/docs/RUNNER_GUIDE.md</code> · Ledger: <code>tier2/docs/VERIFICATION_LEDGER.md</code></li>
<li>Dead ends: <code>tier2/docs/DEAD_ENDS_REGISTER.md</code> · Assumptions: <code>tier2/docs/ASSUMPTIONS_REGISTER.md</code></li>
<li>All raw evidence: <code>tier2/evidence/</code> · Harnesses/GDB scripts: <code>tier2/scripts/</code></li>
<li>Build &amp; run: <code>DEBUG=1 ./scripts/run_qemu.sh</code>, <code>gdb -batch -q -x scripts/&lt;S&gt;.py android/artifacts/vmlinux</code></li>
</ul>
</body></html>"""
    OUT_PDF.parent.mkdir(parents=True, exist_ok=True)
    try:
        weasyprint.HTML(string=html_doc).write_pdf(str(OUT_PDF))
        print(f"[+] PDF backup written: {OUT_PDF}")
    except Exception as exc:
        print(f"[!] PDF backup failed: {exc}", file=sys.stderr)


def main():
    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"[+] Deck saved: {OUT_PPTX} ({_slide_no[0]} slides)")
    build_pdf_backup()


if __name__ == "__main__":
    main()
